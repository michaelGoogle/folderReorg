"""
Full-text extraction for the knowledge base.

Differences from src/phase1_extract.py:
  · No 4000-char cap — we want the FULL document.
  · PDFs with no extractable text are OCR'd with Tesseract
    (configurable languages; default DEU+ENG).
  · Stand-alone images (.jpg, .png, …) are OCR'd too.
  · Returns one string per file (joined from pages / cells / paragraphs).
"""

from __future__ import annotations

from pathlib import Path

from src.config import MAX_TEXT_CHARS  # not used here, just to stay consistent
from kb.config import (
    MAX_FILE_SIZE_BYTES,
    OCR_DPI,
    OCR_ENABLED,
    OCR_IMAGE_EXTS,
    OCR_LANGS,
    TEXT_EXTS,
)


# --- MuPDF stderr taming ---------------------------------------------------
# Without this, PyMuPDF lets its C library scream lines like
#   "MuPDF error: format error: corrupt object stream (47 0 R)"
# straight onto our process stderr, with no filename context — which is
# useless when we're streaming through thousands of PDFs in one run.
# We disable that, route everything into PyMuPDF's in-memory warning buffer,
# and drain it per-file (see `_drain_mupdf` below) so the offending filename
# is right next to the error.
def _silence_mupdf_default_stderr() -> None:
    try:
        import fitz
        # Stop printing to stderr; warnings now accumulate in fitz.TOOLS.mupdf_warnings.
        if hasattr(fitz, "TOOLS"):
            try:
                fitz.TOOLS.mupdf_display_errors(False)
            except Exception:
                pass
            try:
                fitz.TOOLS.mupdf_display_warnings(False)
            except Exception:
                pass
    except Exception:
        # PyMuPDF not importable yet — _extract_pdf will raise its own error.
        pass

_silence_mupdf_default_stderr()


# Patterns for recoverable mupdf chatter that we suppress by default. These
# are emitted by the C library while it auto-repairs minor PDF defects —
# extraction succeeds, the warnings are noise. The list is matched as
# substrings against each line (case-insensitive). Override by setting
# KB_MUPDF_VERBOSE=1 in the environment to see everything (useful when
# debugging a specific bad PDF).
import os as _os
import re as _re

_MUPDF_HARMLESS_PATTERNS = [
    # XREF / trailer auto-repair (mupdf rebuilds the table and proceeds)
    "broken xref subsection, proceeding anyway",
    "incorrect number of xref entries in trailer, repairing",
    "trying to repair broken xref",
    "repairing pdf document",
    "trailer size is off-by-one. ignoring",
    # Stream-length mismatch — mupdf computes the true length and uses that
    "pdf stream length incorrect",
    "line feed missing after stream begin marker",
    # Image colorspace mismatches inside PDFs (visual only, doesn't affect text)
    "jpx numcomps",
    # Drawing-state warnings (curveto/lineto/closepath without current point) —
    # affect rendering only, not text extraction
    "curveto with no current point",
    "lineto with no current point",
    "closepath with no current point",
    # Outline / metadata repairs
    "repaired broken tree structure in outline",
    "bogus font ascent/descent",
    # CMap warnings — most are recoverable; if extraction is empty we'll see
    # that downstream via res.text being short anyway
    "cmap is missing codespace range",
    # Decompression chatter that mupdf says it's ignoring
    "ignoring zlib error",
    # Repeat-count line that fitz sometimes appends
    "... repeated ",
]

# Patterns ALWAYS shown — these are real signal:
#   · password / encryption errors      → file likely unindexable
#   · "format error: corrupt object"    → real damage; extraction may be partial
#   · "non-page object in page tree"    → page-tree damage; data loss probable
#   · "cannot load content stream"      → text loss
#   · "actualtext with no position"     → text loss (caller may want to know)
#   · plain "syntax error"              → real corruption attempt
#   · "cannot load object"              → real corruption
#   · "page tree load failed"           → real corruption
# We don't need an explicit allowlist — anything not matched by the harmless
# list above falls through and gets printed.

_MUPDF_VERBOSE = _os.environ.get("KB_MUPDF_VERBOSE", "0") not in ("", "0", "false", "no")


def _is_harmless_mupdf(line: str) -> bool:
    if _MUPDF_VERBOSE:
        return False
    low = line.lower()
    return any(p in low for p in _MUPDF_HARMLESS_PATTERNS)


def _drain_mupdf(path: Path, op: str) -> None:
    """
    Pull any MuPDF warnings/errors that accumulated during the most recent
    fitz operation on `path`. Re-emit only the actionable ones (corruption,
    password issues, page-tree failures, …) tagged with the filename.

    Recoverable chatter (auto-repaired xref, stream length mismatches,
    drawing-state warnings, JPX colorspace mismatches, etc.) is suppressed.
    Set the env var KB_MUPDF_VERBOSE=1 to see everything when debugging.

    Uses tqdm.write when tqdm is active so the progress bar isn't garbled.
    """
    try:
        import fitz
        msg = fitz.TOOLS.mupdf_warnings(reset=True) if hasattr(fitz, "TOOLS") else ""
    except Exception:
        return
    if not msg:
        return
    try:
        from tqdm import tqdm as _tqdm
        emit = _tqdm.write
    except Exception:
        import sys
        def emit(s: str) -> None:
            print(s, file=sys.stderr, flush=True)

    seen: set[str] = set()
    suppressed = 0
    for line in msg.splitlines():
        line = line.strip()
        if not line or line in seen:
            continue
        seen.add(line)
        if _is_harmless_mupdf(line):
            suppressed += 1
            continue
        emit(f"  MuPDF [{op}] {path.name}: {line}")
    # Only mention suppression when there were a LOT of warnings — keeps
    # quick scans clean while still hinting that mupdf did some work.
    if suppressed >= 10:
        emit(f"  MuPDF [{op}] {path.name}: ({suppressed} recoverable "
             f"warnings auto-repaired; set KB_MUPDF_VERBOSE=1 to see them)")


class ExtractResult:
    __slots__ = ("text", "status", "pages", "ocr_used")

    def __init__(self, text: str, status: str, pages: int = 0, ocr_used: bool = False):
        self.text = text
        self.status = status       # "ok", "password", "corrupt", "too_large", "unsupported", "empty"
        self.pages = pages
        self.ocr_used = ocr_used


# ------------------------------------------------------------------ PDFs ---


def _pdf_text_first_pass(path: Path) -> tuple[str, int, bool]:
    """Native PyMuPDF text extraction. Returns (text, num_pages, encrypted)."""
    import fitz  # PyMuPDF
    try:
        with fitz.open(path) as doc:
            if doc.is_encrypted:
                return "", len(doc), True
            parts = []
            for page in doc:
                parts.append(page.get_text("text"))
            return "\n".join(parts).strip(), len(doc), False
    finally:
        _drain_mupdf(path, "text")


def _pdf_ocr(path: Path) -> str:
    """Render each page to an image and OCR with Tesseract."""
    import fitz
    import pytesseract
    from PIL import Image
    import io
    chunks: list[str] = []
    try:
        with fitz.open(path) as doc:
            for page in doc:
                pix = page.get_pixmap(dpi=OCR_DPI, alpha=False)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = pytesseract.image_to_string(img, lang=OCR_LANGS)
                if text.strip():
                    chunks.append(text)
    finally:
        _drain_mupdf(path, "ocr")
    return "\n\n".join(chunks).strip()


def _extract_pdf(path: Path) -> ExtractResult:
    try:
        text, pages, encrypted = _pdf_text_first_pass(path)
    except Exception as e:
        # Make sure the filename appears in the log even when fitz.open
        # itself raises (e.g. truly broken PDF, password-protected with
        # AES, etc.) — _drain_mupdf already ran inside the helper.
        try:
            from tqdm import tqdm as _tqdm
            _tqdm.write(f"  MuPDF [open-fail] {path.name}: {e}")
        except Exception:
            import sys
            print(f"  MuPDF [open-fail] {path.name}: {e}",
                  file=sys.stderr, flush=True)
        return ExtractResult("", "corrupt")
    if encrypted:
        return ExtractResult("", "password", pages=pages)
    # Heuristic: if native extraction yielded < ~50 chars per page, assume
    # image-only PDF and OCR it.
    avg = len(text) / max(pages, 1)
    if avg < 50 and OCR_ENABLED:
        try:
            ocr_text = _pdf_ocr(path)
        except Exception:
            return ExtractResult(text, "ok" if text else "corrupt",
                                 pages=pages, ocr_used=False)
        if len(ocr_text) > len(text):
            return ExtractResult(ocr_text, "ok", pages=pages, ocr_used=True)
    return ExtractResult(text, "ok" if text else "empty", pages=pages)


# ------------------------------------------------------------------ Office --


def _extract_docx(path: Path) -> ExtractResult:
    try:
        from docx import Document as Docx
        doc = Docx(path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
        return ExtractResult(text, "ok" if text else "empty")
    except Exception:
        return ExtractResult("", "corrupt")


def _extract_xlsx(path: Path) -> ExtractResult:
    """Modern Excel (.xlsx, .xlsm) via openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c for c in cells):
                    parts.append("\t".join(cells))
        return ExtractResult("\n".join(parts), "ok" if parts else "empty")
    except Exception:
        return ExtractResult("", "corrupt")


def _extract_xls(path: Path) -> ExtractResult:
    """Legacy Excel (.xls) via xlrd. Note: xlrd >= 2.0 dropped .xlsx
    support, so this codepath is .xls-only — .xlsx still routes to
    _extract_xlsx via the dispatcher."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts: list[str] = []
        for sheet in wb.sheets():
            parts.append(f"# Sheet: {sheet.name}")
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                cells = ["" if v is None or v == "" else str(v) for v in row]
                if any(c for c in cells):
                    parts.append("\t".join(cells))
        return ExtractResult("\n".join(parts), "ok" if parts else "empty")
    except Exception:
        return ExtractResult("", "corrupt")


def _extract_pptx(path: Path) -> ExtractResult:
    """Modern PowerPoint (.pptx) via python-pptx. Walks every slide and
    every text-bearing shape (text frames, tables, notes)."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts: list[str] = []
        for slide_no, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                # Plain text-frame shapes
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in para.runs).strip()
                        if line:
                            slide_parts.append(line)
                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            slide_parts.append("\t".join(cells))
            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"[notes] {notes}")
            if slide_parts:
                parts.append(f"# Slide {slide_no}")
                parts.extend(slide_parts)
        return ExtractResult("\n".join(parts), "ok" if parts else "empty")
    except Exception:
        return ExtractResult("", "corrupt")


def _extract_doc_via_antiword(path: Path) -> ExtractResult:
    """Legacy Word (.doc) via the `antiword` system binary. Returns
    `unsupported` (not corrupt) if antiword isn't installed, so the
    indexer's synthetic-context fallback kicks in cleanly."""
    import shutil as _shutil
    import subprocess as _sp
    if not _shutil.which("antiword"):
        return ExtractResult("", "unsupported")
    try:
        r = _sp.run(["antiword", "-w", "0", str(path)],
                    capture_output=True, timeout=30)
        if r.returncode != 0:
            return ExtractResult("", "corrupt")
        text = r.stdout.decode("utf-8", errors="replace").strip()
        return ExtractResult(text, "ok" if text else "empty")
    except _sp.TimeoutExpired:
        return ExtractResult("", "corrupt")
    except Exception:
        return ExtractResult("", "corrupt")


def _extract_plain(path: Path) -> ExtractResult:
    try:
        text = path.read_text(errors="ignore")
        return ExtractResult(text, "ok" if text.strip() else "empty")
    except Exception:
        return ExtractResult("", "corrupt")


# ------------------------------------------------------------------ Images -


def _ocr_image(path: Path) -> ExtractResult:
    if not OCR_ENABLED:
        return ExtractResult("", "unsupported")
    try:
        import pytesseract
        from PIL import Image
        with Image.open(path) as img:
            text = pytesseract.image_to_string(img, lang=OCR_LANGS)
        return ExtractResult(text.strip(), "ok" if text.strip() else "empty",
                             ocr_used=True)
    except Exception:
        return ExtractResult("", "corrupt")


# ------------------------------------------------------------------ dispatch


def extract(path: Path) -> ExtractResult:
    """Extract full text from `path`. Never raises; returns status on failure.

    Files whose extension we don't have an extractor for (or where the
    extractor can't yield text — e.g. .ppt without LibreOffice, encrypted
    PDFs, image-only PDFs OCR'd to nothing) return status != "ok". The
    indexer then falls back to indexing a synthetic context document built
    from the filename + folder hierarchy, so the file is still findable
    via filename / folder semantics in chat search.
    """
    if not path.exists() or not path.is_file():
        return ExtractResult("", "corrupt")
    try:
        size = path.stat().st_size
    except OSError:
        return ExtractResult("", "corrupt")
    if size > MAX_FILE_SIZE_BYTES:
        return ExtractResult("", "too_large")
    ext = path.suffix.lower()
    # PDFs (with PyMuPDF + Tesseract OCR fallback for image-only)
    if ext == ".pdf":
        return _extract_pdf(path)
    # Word
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".doc":
        # Requires the `antiword` system binary; falls through to
        # "unsupported" (→ synthetic context) if not installed.
        return _extract_doc_via_antiword(path)
    # Excel
    if ext in (".xlsx", ".xlsm"):
        return _extract_xlsx(path)
    if ext == ".xls":
        return _extract_xls(path)
    # PowerPoint
    if ext == ".pptx":
        return _extract_pptx(path)
    # .ppt has no clean lightweight extractor — needs LibreOffice
    # conversion. Returns "unsupported" so synthetic context kicks in.
    # Plain text variants
    if ext in (".txt", ".md", ".csv", ".rtf"):
        return _extract_plain(path)
    # Standalone images via OCR
    if ext in OCR_IMAGE_EXTS:
        return _ocr_image(path)
    return ExtractResult("", "unsupported")
